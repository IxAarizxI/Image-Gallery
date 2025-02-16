from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Add a slide with a title and content layout
slide_layout = prs.slide_layouts[5]  # Blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Define positions and sizes for shapes
left = Inches(1)
top = Inches(1)
width = Inches(2)
height = Inches(1)

# Add shapes for each module
modules = [
    ("Main Application", Inches(3), Inches(0.5)),
    ("Image Gallery", Inches(1), Inches(2)),
    ("Edit Module", Inches(4), Inches(2)),
    ("Resize Module", Inches(1), Inches(3.5)),
    ("Crop Module", Inches(4), Inches(3.5)),
    ("Slideshow Module", Inches(2.5), Inches(5)),
]

# Add shapes and text
for module_name, x, y in modules:
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape
        x, y, width, height
    )
    shape.text = module_name
    shape.text_frame.paragraphs[0].alignment = 1  # Center align text

# Add connectors (arrows) between shapes
connector = slide.shapes.add_connector(
    3,  # Straight line connector
    Inches(3.5), Inches(1.5),  # Start point (Main Application bottom)
    Inches(3.5), Inches(2)     # End point (Image Gallery top)
)

connector = slide.shapes.add_connector(
    3,  # Straight line connector
    Inches(2), Inches(2.5),  # Start point (Image Gallery bottom)
    Inches(2), Inches(3.5)  # End point (Resize Module top)
)

connector = slide.shapes.add_connector(
    3,  # Straight line connector
    Inches(4), Inches(2.5),  # Start point (Image Gallery bottom)
    Inches(4), Inches(3.5)   # End point (Crop Module top)
)

connector = slide.shapes.add_connector(
    3,  # Straight line connector
    Inches(2.5), Inches(4.5),  # Start point (Resize Module bottom)
    Inches(2.5), Inches(5)   # End point (Slideshow Module top)
)

# Save the presentation
prs.save("Image_Gallery_Block_Diagram.pptx")

print("Block diagram created successfully!")